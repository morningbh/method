"""File processor service — upload persistence + pdf/docx text extraction.

Task 3.1 (Issue #2 / M3). See
``docs/design/issue-2-task-3.1-file-processor.md`` for the full contract.

Responsibilities
----------------
- ``validate_upload_limits`` — reject request-level limit violations using
  declared filename/size (no content read needed).
- ``save_and_extract`` — MIME-sniff, persist original to disk under a
  per-request directory (``{upload_dir}/{request_id}/``), and, for pdf/docx,
  offload text extraction to the default thread-pool executor with a 10s
  timeout. Write extracted text to a sibling ``.extracted.md`` so the
  downstream claude subprocess only needs the ``Read`` tool on plain
  markdown (HARNESS §3).
- ``cleanup_request`` — recursively remove a request directory; idempotent.

Security notes
--------------
- ``request_id`` is a regex-validated ULID before use as a directory name
  (path-traversal defense, design §2).
- ``original_name`` is never joined into a filesystem path; on-disk name is
  ``{uuid4.hex}{ext}``.
- Extraction failures are best-effort: logged and reported via
  ``SavedFile.extraction_ok=False``, never raised. The original file is
  always preserved.
"""
from __future__ import annotations

import asyncio
import logging
import re
import shutil
import uuid
from dataclasses import dataclass
from pathlib import Path

import magic
from fastapi import HTTPException, UploadFile

from app import config as _config

logger = logging.getLogger("method.file_processor")

# Crockford base32 (26 chars, no I/L/O/U). Matches the canonical ULID spec.
_ULID_RE = re.compile(r"^[0-9A-HJKMNP-TV-Z]{26}$")

_ALLOWED_EXTS = frozenset(
    {
        # Text
        ".md", ".txt",
        # Office
        ".pdf", ".docx", ".pptx", ".xlsx",
        # Images (Claude Read handles visually, no server-side extraction)
        ".png", ".jpg", ".jpeg", ".webp", ".gif",
    }
)
_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".webp", ".gif"})

_MAX_FILES = 20
_MAX_FILE_BYTES = 50 * 1024 * 1024
_MAX_TOTAL_BYTES = 100 * 1024 * 1024
_EXTRACTION_TIMEOUT = 10.0
_SNIFF_BYTES = 2048

# xlsx extraction: truncate each sheet to this many rows. Protects the prompt
# token budget — a 10 MB spreadsheet can otherwise produce multi-MB of text.
_XLSX_MAX_ROWS_PER_SHEET = 500

# Per-extension accepted sniffed MIMEs (design §3). libmagic often reports
# ``application/zip`` for OOXML formats (docx/pptx/xlsx) because the outer
# container is a zip — accept both that and the fully-qualified OOXML type.
_ACCEPTED_MIMES: dict[str, frozenset[str]] = {
    ".md": frozenset({"text/plain", "text/markdown", "text/x-markdown"}),
    ".txt": frozenset({"text/plain"}),
    ".pdf": frozenset({"application/pdf"}),
    ".docx": frozenset(
        {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/zip",
        }
    ),
    ".pptx": frozenset(
        {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/zip",
        }
    ),
    ".xlsx": frozenset(
        {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/zip",
        }
    ),
    ".png": frozenset({"image/png"}),
    ".jpg": frozenset({"image/jpeg"}),
    ".jpeg": frozenset({"image/jpeg"}),
    ".webp": frozenset({"image/webp"}),
    ".gif": frozenset({"image/gif"}),
}

# Module-level singleton — magic.Magic is not cheap to construct and is
# thread-safe for read-only sniffing via from_buffer.
_MAGIC = magic.Magic(mime=True)


@dataclass(frozen=True)
class SavedFile:
    """Result of ``save_and_extract``. All paths are absolute (HARNESS §2)."""

    stored_path: Path
    extracted_path: Path | None
    size_bytes: int
    mime_type: str
    extraction_ok: bool


class LimitExceededError(HTTPException):
    """HTTP 400 with ``detail = {"code": <code>, "message": <str>}``.

    Six known codes: ``files_too_many``, ``file_too_large``,
    ``total_too_large``, ``unsupported_type``, ``empty_file``,
    ``mime_mismatch``.
    """

    def __init__(self, code: str, message: str) -> None:
        super().__init__(
            status_code=400, detail={"code": code, "message": message}
        )


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


async def validate_upload_limits(files: list[UploadFile]) -> None:
    """Raise ``LimitExceededError`` on request-level limit violations.

    Uses declared filename + size only (design §2). MIME sniff happens in
    ``save_and_extract`` because that's where the bytes are available.
    """
    if len(files) > _MAX_FILES:
        raise LimitExceededError(
            "files_too_many",
            f"{len(files)} files exceeds the {_MAX_FILES}-per-request limit",
        )

    total = 0
    for f in files:
        name = f.filename or ""
        ext = Path(name).suffix.lower()
        if ext not in _ALLOWED_EXTS:
            raise LimitExceededError(
                "unsupported_type",
                f"{name!r}: extension {ext!r} not in {sorted(_ALLOWED_EXTS)}",
            )
        size = f.size or 0
        if size == 0:
            raise LimitExceededError("empty_file", f"{name!r} is empty (0 bytes)")
        if size > _MAX_FILE_BYTES:
            raise LimitExceededError(
                "file_too_large",
                f"{name!r} is {size} bytes, exceeds {_MAX_FILE_BYTES}",
            )
        total += size

    if total > _MAX_TOTAL_BYTES:
        raise LimitExceededError(
            "total_too_large",
            f"total {total} bytes exceeds {_MAX_TOTAL_BYTES}",
        )


async def save_and_extract(
    request_id: str,
    original_name: str,
    content: bytes,
) -> SavedFile:
    """Persist ``content`` and, for pdf/docx, extract text to a sibling file.

    See module docstring and ``docs/design/issue-2-task-3.1-file-processor.md``
    §2 for the full contract.
    """
    # 1. Trust-boundary check — request_id is used as a directory component,
    #    a malformed value would allow path traversal.
    if not _ULID_RE.fullmatch(request_id):
        raise ValueError(f"invalid request_id: {request_id!r}")

    ext = Path(original_name).suffix.lower()
    if ext not in _ALLOWED_EXTS:
        raise LimitExceededError(
            "unsupported_type",
            f"{original_name!r}: extension {ext!r} not in {sorted(_ALLOWED_EXTS)}",
        )

    # 2. MIME sniff against the actual bytes (spoof defense).
    sniffed = _MAGIC.from_buffer(content[:_SNIFF_BYTES])
    accepted = _ACCEPTED_MIMES[ext]
    if sniffed not in accepted:
        raise LimitExceededError(
            "mime_mismatch",
            f"{original_name!r}: sniffed {sniffed!r} not in {sorted(accepted)}",
        )

    # 3. mkdir -p {upload_dir}/{request_id}
    upload_root = Path(_config.settings.upload_dir).resolve()
    req_dir = upload_root / request_id
    req_dir.mkdir(parents=True, exist_ok=True)

    # 4. Write content to {uuid4.hex}{ext}
    stem = uuid.uuid4().hex
    stored_path = (req_dir / f"{stem}{ext}").resolve()
    stored_path.write_bytes(content)

    # 5. Extraction with timeout + broad failure catch.
    #    md / txt / images: nothing to extract — stored as-is. For images,
    #    Claude's Read tool handles visual reading at planning time.
    #    pdf / docx / pptx / xlsx: extract to a sibling .extracted.md so the
    #    downstream claude subprocess only needs plain text (HARNESS §3).
    extracted_path: Path | None = None
    extraction_ok = True

    _EXTRACTORS = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
    }
    if ext in _EXTRACTORS:
        text = await _run_extractor(_EXTRACTORS[ext], stored_path)
        if text:
            extracted_path = (req_dir / f"{stem}.extracted.md").resolve()
            extracted_path.write_text(text, encoding="utf-8")
            extraction_ok = True
        else:
            extracted_path = None
            extraction_ok = False

    return SavedFile(
        stored_path=stored_path,
        extracted_path=extracted_path,
        size_bytes=len(content),
        mime_type=sniffed,
        extraction_ok=extraction_ok,
    )


async def cleanup_request(request_id: str) -> None:
    """Recursively delete ``{upload_dir}/{request_id}/``. Idempotent."""
    if not _ULID_RE.fullmatch(request_id):
        raise ValueError(f"invalid request_id: {request_id!r}")
    req_dir = Path(_config.settings.upload_dir).resolve() / request_id
    if req_dir.exists():
        shutil.rmtree(req_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# Private helpers
# ---------------------------------------------------------------------------


async def _run_extractor(extractor, path: Path) -> str | None:
    """Run a blocking extractor in the default executor with a timeout.

    Returns the non-empty extracted text on success, or ``None`` on any
    failure (timeout, parser error, encrypted source, empty output). Never
    raises to the caller — best-effort, per design §4.
    """
    loop = asyncio.get_running_loop()
    try:
        text = await asyncio.wait_for(
            loop.run_in_executor(None, extractor, path),
            timeout=_EXTRACTION_TIMEOUT,
        )
    except TimeoutError:
        logger.warning(
            "extraction_timeout path=%s timeout=%ss", path, _EXTRACTION_TIMEOUT
        )
        return None
    except Exception as exc:  # noqa: BLE001 — best-effort per design §4
        logger.warning("extraction_failed path=%s err=%r", path, exc)
        return None

    if not text:
        logger.warning("extraction_empty path=%s", path)
        return None
    return text


def _extract_pdf(path: Path) -> str:
    """Extract text from a PDF via pdfplumber. Blocking — call in executor."""
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        return "\n\n".join((page.extract_text() or "") for page in pdf.pages).strip()


def _extract_docx(path: Path) -> str:
    """Extract text from a .docx via python-docx. Blocking — call in executor."""
    from docx import Document

    return "\n".join(p.text for p in Document(path).paragraphs).strip()


def _extract_pptx(path: Path) -> str:
    """Extract text from a .pptx via python-pptx. Blocking — call in executor.

    Emits per-slide sections with shape text + speaker notes. No slide-image
    OCR — that's Claude Read's job if the user needs it.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    out.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append("")
                out.append(f"_Notes:_ {notes}")
        out.append("")
    return "\n".join(out).strip()


def _extract_xlsx(path: Path) -> str:
    """Extract text from a .xlsx via openpyxl (read-only + values-only mode).

    Blocking — call in executor. Each sheet is dumped as tab-separated rows
    up to ``_XLSX_MAX_ROWS_PER_SHEET``; past that, a truncation marker is
    appended. ``read_only=True`` streams rows (low RAM for large workbooks);
    ``data_only=True`` returns cached cell values rather than formula text.
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        out: list[str] = []
        for sheet in wb.worksheets:
            out.append(f"## Sheet: {sheet.title}")
            truncated = False
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= _XLSX_MAX_ROWS_PER_SHEET:
                    truncated = True
                    break
                cells = [
                    "" if v is None else str(v).replace("\t", " ").replace("\n", " ")
                    for v in row
                ]
                # Drop wholly-empty trailing cells for readability.
                while cells and cells[-1] == "":
                    cells.pop()
                if cells:
                    out.append("\t".join(cells))
            if truncated:
                out.append(f"[... 已截断，仅展示前 {_XLSX_MAX_ROWS_PER_SHEET} 行]")
            out.append("")
        return "\n".join(out).strip()
    finally:
        wb.close()
